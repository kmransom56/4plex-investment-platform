"""
Export and Reporting System for 4-plex Investment Platform
Generates professional reports in multiple formats integrating multifamily functionality
"""

import os
import json
import logging
from typing import Dict, List, Optional, Any, Union
from datetime import datetime, timedelta
from pathlib import Path
import uuid
import shutil
import asyncio
from concurrent.futures import ThreadPoolExecutor
import io
import base64

# Excel generation
try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.chart import LineChart, BarChart, Reference
    from openpyxl.utils.dataframe import dataframe_to_rows
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logging.warning("Excel libraries not available. Excel generation disabled.")

# PDF generation
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("PDF libraries not available. PDF generation disabled.")

# PowerPoint generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("PowerPoint libraries not available. PowerPoint generation disabled.")

from models import Property, PropertyAnalysis, FinancialMetrics
from database.connection import DatabaseManager
from financial_calculations import FinancialCalculator

class ExportSystem:
    """
    Comprehensive export and reporting system
    Generates professional reports in multiple formats
    """
    
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self.db_manager = DatabaseManager(config.get('database', {}))
        self.financial_calculator = FinancialCalculator(config.get('financial_assumptions', {}))
        
        # Export configuration
        self.export_dir = config.get('export', {}).get('output_dir', './exports')
        self.max_export_age_days = config.get('export', {}).get('max_export_age_days', 7)
        self.template_dir = config.get('export', {}).get('template_dir', './templates')
        
        # Ensure directories exist
        os.makedirs(self.export_dir, exist_ok=True)
        os.makedirs(self.template_dir, exist_ok=True)
        
        # Brand styling
        self.brand_colors = {
            'primary': '#1f4e79',
            'secondary': '#70ad47',
            'accent': '#ffc000',
            'text': '#333333',
            'light_gray': '#f2f2f2',
            'success': '#28a745',
            'warning': '#ffc107',
            'danger': '#dc3545'
        }
        
        # Thread pool for background processing
        self.executor = ThreadPoolExecutor(max_workers=3)
    
    async def export_property_analysis(
        self,
        property_id: str,
        analysis_id: Optional[str] = None,
        export_formats: List[str] = ['excel', 'pdf', 'json'],
        include_charts: bool = True,
        include_raw_data: bool = False
    ) -> Dict[str, Any]:
        """
        Export comprehensive property analysis
        
        Args:
            property_id: Property to export
            analysis_id: Specific analysis ID (uses latest if None)
            export_formats: List of formats to generate
            include_charts: Include charts in reports
            include_raw_data: Include raw document data
            
        Returns:
            Export results with file paths and metadata
        """
        try:
            # Get property and analysis data
            property_data = self.db_manager.get_property(property_id)
            if not property_data:
                raise ValueError(f"Property not found: {property_id}")
            
            if analysis_id:
                analysis_data = self.db_manager.get_analysis(analysis_id)
            else:
                analysis_data = self.db_manager.get_latest_analysis(property_id)
            
            if not analysis_data:
                raise ValueError(f"No analysis found for property: {property_id}")
            
            # Create export session
            export_id = f"export_{uuid.uuid4().hex[:12]}"
            export_session = {
                'export_id': export_id,
                'property_id': property_id,
                'analysis_id': analysis_data['id'],
                'requested_formats': export_formats,
                'created_at': datetime.now(),
                'status': 'processing',
                'files': []
            }
            
            # Create export directory
            export_path = os.path.join(self.export_dir, export_id)
            os.makedirs(export_path, exist_ok=True)
            
            self.logger.info(f"Starting export {export_id} for property {property_id}")
            
            # Generate reports in parallel
            export_tasks = []
            
            if 'json' in export_formats:
                export_tasks.append(
                    self._export_json_data(property_data, analysis_data, export_path, include_raw_data)
                )
            
            if 'excel' in export_formats and EXCEL_AVAILABLE:
                export_tasks.append(
                    self._export_excel_report(property_data, analysis_data, export_path, include_charts)
                )
            
            if 'pdf' in export_formats and PDF_AVAILABLE:
                export_tasks.append(
                    self._export_pdf_report(property_data, analysis_data, export_path)
                )
            
            if 'powerpoint' in export_formats and PPTX_AVAILABLE:
                export_tasks.append(
                    self._export_powerpoint_report(property_data, analysis_data, export_path)
                )
            
            if 'csv' in export_formats:
                export_tasks.append(
                    self._export_csv_data(property_data, analysis_data, export_path)
                )
            
            # Execute export tasks
            results = await asyncio.gather(*export_tasks, return_exceptions=True)
            
            # Process results
            successful_exports = []
            export_errors = []
            
            for i, result in enumerate(results):
                if isinstance(result, Exception):
                    export_errors.append(str(result))
                    self.logger.error(f"Export task {i} failed: {result}")
                elif result:
                    successful_exports.append(result)
                    export_session['files'].append(result)
            
            # Generate summary report
            if successful_exports:
                summary = self._generate_export_summary(export_session, export_path)
                if summary:
                    export_session['files'].append(summary)
            
            # Update session status
            export_session['status'] = 'completed' if successful_exports else 'failed'
            export_session['completed_at'] = datetime.now()
            export_session['errors'] = export_errors
            
            # Save session metadata
            session_file = os.path.join(export_path, 'export_session.json')
            with open(session_file, 'w') as f:
                json.dump(export_session, f, indent=2, default=str)
            
            self.logger.info(f"Export {export_id} completed: {len(successful_exports)} files generated")
            
            return {
                'export_id': export_id,
                'status': export_session['status'],
                'export_path': export_path,
                'files': export_session['files'],
                'errors': export_errors,
                'total_files': len(successful_exports),
                'created_at': export_session['created_at'].isoformat(),
                'completed_at': export_session['completed_at'].isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting property analysis: {str(e)}")
            raise
    
    async def export_portfolio_summary(
        self,
        filters: Optional[Dict[str, Any]] = None,
        export_formats: List[str] = ['excel', 'pdf'],
        include_property_details: bool = True
    ) -> Dict[str, Any]:
        """
        Export portfolio summary report
        
        Args:
            filters: Optional filters for properties to include
            export_formats: List of formats to generate
            include_property_details: Include individual property details
            
        Returns:
            Export results with file paths
        """
        try:
            # Get portfolio data
            properties = await self.db_manager.get_properties(filters=filters, limit=1000)
            
            if not properties:
                raise ValueError("No properties found matching criteria")
            
            # Get portfolio analytics
            portfolio_summary = await self._calculate_portfolio_analytics(properties)
            
            # Create export session
            export_id = f"portfolio_{uuid.uuid4().hex[:12]}"
            export_path = os.path.join(self.export_dir, export_id)
            os.makedirs(export_path, exist_ok=True)
            
            self.logger.info(f"Exporting portfolio summary: {len(properties)} properties")
            
            # Generate reports
            export_results = []
            
            if 'json' in export_formats:
                json_result = await self._export_portfolio_json(
                    portfolio_summary, properties, export_path, include_property_details
                )
                if json_result:
                    export_results.append(json_result)
            
            if 'excel' in export_formats and EXCEL_AVAILABLE:
                excel_result = await self._export_portfolio_excel(
                    portfolio_summary, properties, export_path, include_property_details
                )
                if excel_result:
                    export_results.append(excel_result)
            
            if 'pdf' in export_formats and PDF_AVAILABLE:
                pdf_result = await self._export_portfolio_pdf(
                    portfolio_summary, properties, export_path
                )
                if pdf_result:
                    export_results.append(pdf_result)
            
            return {
                'export_id': export_id,
                'status': 'completed',
                'export_path': export_path,
                'files': export_results,
                'property_count': len(properties),
                'created_at': datetime.now().isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting portfolio summary: {str(e)}")
            raise
    
    async def export_investment_opportunity_report(
        self,
        min_grade: str = 'B',
        max_properties: int = 50,
        export_format: str = 'pdf'
    ) -> Dict[str, Any]:
        """
        Export investment opportunities report
        
        Args:
            min_grade: Minimum investment grade
            max_properties: Maximum number of properties to include
            export_format: Output format
            
        Returns:
            Export results
        """
        try:
            # Get investment opportunities
            grade_hierarchy = {'A': 4, 'B': 3, 'C': 2, 'D': 1}
            min_grade_value = grade_hierarchy.get(min_grade, 3)
            grade_filters = [grade for grade, value in grade_hierarchy.items() if value >= min_grade_value]
            
            opportunities = await self.db_manager.get_properties(
                filters={
                    'investment_grade': grade_filters,
                    'foreclosure_stage': ['opportunity', 'analyzed']
                },
                sort_by='investment_score DESC',
                limit=max_properties
            )
            
            if not opportunities:
                raise ValueError(f"No investment opportunities found with grade {min_grade} or better")
            
            # Create export
            export_id = f"opportunities_{uuid.uuid4().hex[:12]}"
            export_path = os.path.join(self.export_dir, export_id)
            os.makedirs(export_path, exist_ok=True)
            
            # Generate report based on format
            if export_format == 'pdf' and PDF_AVAILABLE:
                result = await self._export_opportunities_pdf(opportunities, export_path)
            elif export_format == 'excel' and EXCEL_AVAILABLE:
                result = await self._export_opportunities_excel(opportunities, export_path)
            elif export_format == 'json':
                result = await self._export_opportunities_json(opportunities, export_path)
            else:
                raise ValueError(f"Unsupported export format: {export_format}")
            
            return {
                'export_id': export_id,
                'status': 'completed',
                'export_path': export_path,
                'files': [result] if result else [],
                'opportunity_count': len(opportunities),
                'min_grade': min_grade,
                'created_at': datetime.now().isoformat()
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting investment opportunities: {str(e)}")
            raise
    
    # Individual Export Format Methods
    
    async def _export_json_data(
        self,
        property_data: Dict[str, Any],
        analysis_data: Dict[str, Any],
        export_path: str,
        include_raw_data: bool = False
    ) -> Dict[str, Any]:
        """Export comprehensive JSON data"""
        
        try:
            # Compile data
            json_data = {
                'export_metadata': {
                    'generated_at': datetime.now().isoformat(),
                    'export_type': 'property_analysis',
                    'version': '1.0'
                },
                'property': {
                    'basic_info': self._extract_property_basic_info(property_data),
                    'foreclosure_details': self._extract_foreclosure_details(property_data),
                    'location_info': self._extract_location_info(property_data)
                },
                'analysis': {
                    'summary': self._extract_analysis_summary(analysis_data),
                    'financial_metrics': analysis_data.get('financial_metrics', {}),
                    'investment_grade': analysis_data.get('investment_grade'),
                    'viability_score': analysis_data.get('viability_score'),
                    'confidence_score': analysis_data.get('confidence_score')
                },
                'recommendations': analysis_data.get('ai_insights', {}).get('recommendations', []),
                'risk_factors': analysis_data.get('risk_factors', [])
            }
            
            if include_raw_data:
                json_data['raw_data'] = {
                    'property': property_data,
                    'analysis': analysis_data
                }
            
            # Save file
            filename = f"property_analysis_{property_data['id']}.json"
            file_path = os.path.join(export_path, filename)
            
            with open(file_path, 'w') as f:
                json.dump(json_data, f, indent=2, default=str)
            
            return {
                'type': 'json',
                'filename': filename,
                'path': file_path,
                'size_mb': os.path.getsize(file_path) / (1024 * 1024),
                'description': 'Comprehensive property analysis data'
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting JSON: {str(e)}")
            raise
    
    async def _export_excel_report(
        self,
        property_data: Dict[str, Any],
        analysis_data: Dict[str, Any],
        export_path: str,
        include_charts: bool = True
    ) -> Dict[str, Any]:
        """Export detailed Excel report"""
        
        if not EXCEL_AVAILABLE:
            raise ValueError("Excel libraries not available")
        
        try:
            # Create workbook
            wb = Workbook()
            
            # Remove default sheet and create custom sheets
            wb.remove(wb.active)
            
            # Executive Summary sheet
            self._create_executive_summary_sheet(wb, property_data, analysis_data)
            
            # Financial Analysis sheet
            self._create_financial_analysis_sheet(wb, analysis_data, include_charts)
            
            # Investment Projections sheet
            self._create_investment_projections_sheet(wb, analysis_data, include_charts)
            
            # Risk Assessment sheet
            self._create_risk_assessment_sheet(wb, analysis_data)
            
            # Property Details sheet
            self._create_property_details_sheet(wb, property_data)
            
            # Raw Data sheet (if requested)
            if analysis_data.get('processing_metadata', {}).get('include_raw_data', False):
                self._create_raw_data_sheet(wb, property_data, analysis_data)
            
            # Save file
            filename = f"property_analysis_{property_data['id']}.xlsx"
            file_path = os.path.join(export_path, filename)
            wb.save(file_path)
            
            return {
                'type': 'excel',
                'filename': filename,
                'path': file_path,
                'size_mb': os.path.getsize(file_path) / (1024 * 1024),
                'description': 'Detailed Excel analysis with charts and projections'
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting Excel: {str(e)}")
            raise
    
    async def _export_pdf_report(
        self,
        property_data: Dict[str, Any],
        analysis_data: Dict[str, Any],
        export_path: str
    ) -> Dict[str, Any]:
        """Export professional PDF report"""
        
        if not PDF_AVAILABLE:
            raise ValueError("PDF libraries not available")
        
        try:
            filename = f"property_analysis_{property_data['id']}.pdf"
            file_path = os.path.join(export_path, filename)
            
            doc = SimpleDocTemplate(file_path, pagesize=letter)
            styles = getSampleStyleSheet()
            
            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Title'],
                fontSize=24,
                spaceAfter=30,
                textColor=colors.HexColor(self.brand_colors['primary'])
            )
            
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading1'],
                fontSize=16,
                spaceAfter=12,
                textColor=colors.HexColor(self.brand_colors['primary'])
            )
            
            # Build document content
            content = []
            
            # Title page
            content.append(Paragraph("Property Investment Analysis Report", title_style))
            content.append(Spacer(1, 20))
            
            # Property overview
            property_address = f"{property_data.get('address', '')}, {property_data.get('city', '')}, {property_data.get('state', '')}"
            content.append(Paragraph(property_address, styles['Heading2']))
            content.append(Spacer(1, 20))
            
            # Executive summary
            content.append(Paragraph("Executive Summary", heading_style))
            
            investment_grade = analysis_data.get('investment_grade', 'Ungraded')
            viability_score = analysis_data.get('viability_score', 0)
            
            summary_text = f"""
            This property has received an investment grade of <b>{investment_grade}</b> with a viability score of <b>{viability_score:.1%}</b>.
            
            The analysis is based on comprehensive financial modeling, market analysis, and risk assessment.
            """
            
            content.append(Paragraph(summary_text, styles['Normal']))
            content.append(Spacer(1, 20))
            
            # Financial highlights
            content.append(Paragraph("Financial Highlights", heading_style))
            
            financial_metrics = analysis_data.get('financial_metrics', {})
            highlights_data = [
                ['Metric', 'Value'],
                ['Cap Rate', f"{financial_metrics.get('cap_rate', 0):.2%}"],
                ['IRR', f"{financial_metrics.get('irr', 0):.2%}"],
                ['Cash-on-Cash Return', f"{financial_metrics.get('cash_on_cash_return', 0):.2%}"],
                ['NOI', f"${financial_metrics.get('noi', 0):,.0f}"],
                ['DSCR', f"{financial_metrics.get('dscr', 0):.2f}"]
            ]
            
            highlights_table = Table(highlights_data, colWidths=[3*inch, 2*inch])
            highlights_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.brand_colors['primary'])),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            content.append(highlights_table)
            content.append(Spacer(1, 20))
            
            # Investment recommendations
            content.append(Paragraph("Investment Recommendations", heading_style))
            
            recommendations = analysis_data.get('ai_insights', {}).get('recommendations', [])
            if recommendations:
                for i, rec in enumerate(recommendations[:5], 1):
                    content.append(Paragraph(f"{i}. {rec}", styles['Normal']))
            else:
                content.append(Paragraph("No specific recommendations available.", styles['Normal']))
            
            content.append(Spacer(1, 20))
            
            # Risk factors
            content.append(Paragraph("Risk Assessment", heading_style))
            
            risk_factors = analysis_data.get('risk_factors', [])
            if risk_factors:
                for i, risk in enumerate(risk_factors[:5], 1):
                    content.append(Paragraph(f"{i}. {risk}", styles['Normal']))
            else:
                content.append(Paragraph("No significant risk factors identified.", styles['Normal']))
            
            # Build PDF
            doc.build(content)
            
            return {
                'type': 'pdf',
                'filename': filename,
                'path': file_path,
                'size_mb': os.path.getsize(file_path) / (1024 * 1024),
                'description': 'Professional PDF executive summary'
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting PDF: {str(e)}")
            raise
    
    async def _export_powerpoint_report(
        self,
        property_data: Dict[str, Any],
        analysis_data: Dict[str, Any],
        export_path: str
    ) -> Dict[str, Any]:
        """Export PowerPoint investment pitch deck"""
        
        if not PPTX_AVAILABLE:
            raise ValueError("PowerPoint libraries not available")
        
        try:
            # Create presentation
            prs = Presentation()
            
            # Slide 1: Title slide
            title_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Investment Opportunity"
            property_address = f"{property_data.get('address', '')}, {property_data.get('city', '')}"
            subtitle.text = property_address
            
            # Slide 2: Property overview
            bullet_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Property Overview"
            
            tf = content.text_frame
            tf.text = f"Property Type: {property_data.get('property_type', 'Multifamily')}"
            
            p = tf.add_paragraph()
            p.text = f"Total Units: {property_data.get('units', 'TBD')}"
            
            p = tf.add_paragraph()
            p.text = f"Foreclosure Stage: {property_data.get('foreclosure_stage', 'Unknown').title()}"
            
            p = tf.add_paragraph()
            p.text = f"Investment Grade: {analysis_data.get('investment_grade', 'Ungraded')}"
            
            # Slide 3: Financial highlights
            slide = prs.slides.add_slide(bullet_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Financial Highlights"
            
            financial_metrics = analysis_data.get('financial_metrics', {})
            
            tf = content.text_frame
            tf.text = f"Cap Rate: {financial_metrics.get('cap_rate', 0):.2%}"
            
            p = tf.add_paragraph()
            p.text = f"IRR: {financial_metrics.get('irr', 0):.2%}"
            
            p = tf.add_paragraph()
            p.text = f"Cash-on-Cash Return: {financial_metrics.get('cash_on_cash_return', 0):.2%}"
            
            p = tf.add_paragraph()
            p.text = f"NOI: ${financial_metrics.get('noi', 0):,.0f}"
            
            # Slide 4: Investment thesis
            slide = prs.slides.add_slide(bullet_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Investment Thesis"
            
            recommendations = analysis_data.get('ai_insights', {}).get('recommendations', [])
            
            tf = content.text_frame
            if recommendations:
                tf.text = recommendations[0] if recommendations else "Strong investment fundamentals"
                for rec in recommendations[1:4]:
                    p = tf.add_paragraph()
                    p.text = rec
            else:
                tf.text = "Comprehensive analysis supports investment viability"
            
            # Slide 5: Risk considerations
            slide = prs.slides.add_slide(bullet_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Risk Considerations"
            
            risk_factors = analysis_data.get('risk_factors', [])
            
            tf = content.text_frame
            if risk_factors:
                tf.text = risk_factors[0] if risk_factors else "Standard real estate investment risks"
                for risk in risk_factors[1:4]:
                    p = tf.add_paragraph()
                    p.text = risk
            else:
                tf.text = "Low to moderate risk profile based on analysis"
            
            # Save file
            filename = f"investment_pitch_{property_data['id']}.pptx"
            file_path = os.path.join(export_path, filename)
            prs.save(file_path)
            
            return {
                'type': 'powerpoint',
                'filename': filename,
                'path': file_path,
                'size_mb': os.path.getsize(file_path) / (1024 * 1024),
                'description': 'Investment pitch deck presentation'
            }
            
        except Exception as e:
            self.logger.error(f"Error exporting PowerPoint: {str(e)}")
            raise
    
    async def _export_csv_data(
        self,
        property_data: Dict[str, Any],
        analysis_data: Dict[str, Any],
        export_path: str
    ) -> Dict[str, Any]:
        """Export CSV data files"""
        
        try:
            # Create financial metrics CSV
            financial_metrics = analysis_data.get('financial_metrics', {})
            
            metrics_data = []
            for key, value in financial_metrics.items():
                metrics_data.append({
                    'metric': key.replace('_', ' ').title(),
                    'value': value,
                    'property_id': property_data['id']
                })
            
            if metrics_data:
                df = pd.DataFrame(metrics_data)
                csv_filename = f"financial_metrics_{property_data['id']}.csv"
                csv_path = os.path.join(export_path, csv_filename)
                df.to_csv(csv_path, index=False)
                
                return {
                    'type': 'csv',
                    'filename': csv_filename,
                    'path': csv_path,
                    'size_mb': os.path.getsize(csv_path) / (1024 * 1024),
                    'description': 'Financial metrics data'
                }
            
            return None
            
        except Exception as e:
            self.logger.error(f"Error exporting CSV: {str(e)}")
            raise
    
    # Portfolio Export Methods
    
    async def _calculate_portfolio_analytics(self, properties: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Calculate comprehensive portfolio analytics"""
        
        try:
            analytics = {
                'summary': {
                    'total_properties': len(properties),
                    'by_stage': {},
                    'by_grade': {},
                    'by_state': {}
                },
                'financial_metrics': {
                    'total_estimated_value': 0,
                    'total_outstanding_debt': 0,
                    'average_investment_score': 0,
                    'grade_distribution': {}
                },
                'opportunities': {
                    'grade_a_count': 0,
                    'grade_b_count': 0,
                    'high_score_properties': []
                },
                'geographic_distribution': {},
                'analysis_date': datetime.now().isoformat()
            }
            
            # Process each property
            investment_scores = []
            for prop in properties:
                # Stage distribution
                stage = prop.get('foreclosure_stage', 'unknown')
                analytics['summary']['by_stage'][stage] = analytics['summary']['by_stage'].get(stage, 0) + 1
                
                # Grade distribution
                grade = prop.get('investment_grade', 'Ungraded')
                analytics['summary']['by_grade'][grade] = analytics['summary']['by_grade'].get(grade, 0) + 1
                analytics['financial_metrics']['grade_distribution'][grade] = analytics['financial_metrics']['grade_distribution'].get(grade, 0) + 1
                
                # Geographic distribution
                state = prop.get('state', 'Unknown')
                analytics['summary']['by_state'][state] = analytics['summary']['by_state'].get(state, 0) + 1
                analytics['geographic_distribution'][state] = analytics['geographic_distribution'].get(state, 0) + 1
                
                # Financial aggregation
                analytics['financial_metrics']['total_estimated_value'] += prop.get('estimated_value', 0)
                analytics['financial_metrics']['total_outstanding_debt'] += prop.get('outstanding_debt', 0)
                
                # Investment scores
                score = prop.get('investment_score', 0)
                if score > 0:
                    investment_scores.append(score)
                    
                    # High-scoring properties
                    if score > 0.7:
                        analytics['opportunities']['high_score_properties'].append({
                            'id': prop['id'],
                            'address': prop.get('address', ''),
                            'score': score,
                            'grade': grade
                        })
                
                # Opportunity counts
                if grade == 'A':
                    analytics['opportunities']['grade_a_count'] += 1
                elif grade == 'B':
                    analytics['opportunities']['grade_b_count'] += 1
            
            # Calculate averages
            if investment_scores:
                analytics['financial_metrics']['average_investment_score'] = sum(investment_scores) / len(investment_scores)
            
            # Sort high-scoring properties
            analytics['opportunities']['high_score_properties'].sort(key=lambda x: x['score'], reverse=True)
            analytics['opportunities']['high_score_properties'] = analytics['opportunities']['high_score_properties'][:10]
            
            return analytics
            
        except Exception as e:
            self.logger.error(f"Error calculating portfolio analytics: {str(e)}")
            raise
    
    # Excel Sheet Creation Methods
    
    def _create_executive_summary_sheet(self, wb: Workbook, property_data: Dict[str, Any], analysis_data: Dict[str, Any]):
        """Create executive summary Excel sheet"""
        
        ws = wb.create_sheet("Executive Summary")
        
        # Title
        ws['A1'] = "Property Investment Analysis - Executive Summary"
        ws['A1'].font = Font(size=16, bold=True, color=self.brand_colors['primary'])
        ws.merge_cells('A1:F1')
        
        # Property information section
        row = 3
        ws[f'A{row}'] = "Property Information"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        property_info = [
            ('Address:', f"{property_data.get('address', '')}, {property_data.get('city', '')}, {property_data.get('state', '')}"),
            ('Property Type:', property_data.get('property_type', 'Multifamily')),
            ('Foreclosure Stage:', property_data.get('foreclosure_stage', '').title()),
            ('Discovery Date:', str(property_data.get('discovery_date', ''))[:10] if property_data.get('discovery_date') else 'N/A'),
            ('Last Updated:', str(property_data.get('last_modified', ''))[:10] if property_data.get('last_modified') else 'N/A')
        ]
        
        for label, value in property_info:
            ws[f'A{row}'] = label
            ws[f'B{row}'] = value
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
        
        # Investment summary section
        row += 1
        ws[f'A{row}'] = "Investment Summary"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        investment_grade = analysis_data.get('investment_grade', 'Ungraded')
        viability_score = analysis_data.get('viability_score', 0)
        confidence_score = analysis_data.get('confidence_score', 0)
        
        summary_info = [
            ('Investment Grade:', investment_grade),
            ('Viability Score:', f"{viability_score:.1%}"),
            ('Confidence Score:', f"{confidence_score:.1%}"),
            ('Analysis Date:', str(analysis_data.get('analysis_date', ''))[:10] if analysis_data.get('analysis_date') else 'N/A')
        ]
        
        for label, value in summary_info:
            ws[f'A{row}'] = label
            ws[f'B{row}'] = value
            ws[f'A{row}'].font = Font(bold=True)
            
            # Color-code investment grade
            if label == 'Investment Grade:':
                if value == 'A':
                    ws[f'B{row}'].font = Font(color=self.brand_colors['success'], bold=True)
                elif value == 'B':
                    ws[f'B{row}'].font = Font(color=self.brand_colors['secondary'], bold=True)
                elif value in ['C', 'D']:
                    ws[f'B{row}'].font = Font(color=self.brand_colors['warning'], bold=True)
            
            row += 1
        
        # Financial metrics section
        row += 1
        ws[f'A{row}'] = "Key Financial Metrics"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        financial_metrics = analysis_data.get('financial_metrics', {})
        
        # Headers
        ws[f'A{row}'] = "Metric"
        ws[f'B{row}'] = "Value"
        ws[f'A{row}'].font = Font(bold=True)
        ws[f'B{row}'].font = Font(bold=True)
        row += 1
        
        metrics_info = [
            ('Cap Rate', f"{financial_metrics.get('cap_rate', 0):.2%}"),
            ('IRR', f"{financial_metrics.get('irr', 0):.2%}"),
            ('Cash-on-Cash Return', f"{financial_metrics.get('cash_on_cash_return', 0):.2%}"),
            ('NOI', f"${financial_metrics.get('noi', 0):,.0f}"),
            ('DSCR', f"{financial_metrics.get('dscr', 0):.2f}"),
            ('LTV Ratio', f"{financial_metrics.get('ltv_ratio', 0):.2%}")
        ]
        
        for metric, value in metrics_info:
            ws[f'A{row}'] = metric
            ws[f'B{row}'] = value
            row += 1
        
        # Auto-adjust column widths
        for col in ['A', 'B', 'C', 'D', 'E', 'F']:
            ws.column_dimensions[col].auto_size = True
    
    def _create_financial_analysis_sheet(self, wb: Workbook, analysis_data: Dict[str, Any], include_charts: bool = True):
        """Create detailed financial analysis sheet"""
        
        ws = wb.create_sheet("Financial Analysis")
        
        # Title
        ws['A1'] = "Detailed Financial Analysis"
        ws['A1'].font = Font(size=16, bold=True, color=self.brand_colors['primary'])
        ws.merge_cells('A1:H1')
        
        financial_metrics = analysis_data.get('financial_metrics', {})
        
        # Income analysis
        row = 3
        ws[f'A{row}'] = "Income Analysis"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        income_data = [
            ('Gross Income', financial_metrics.get('gross_income', 0)),
            ('Vacancy Allowance', financial_metrics.get('vacancy_allowance', 0)),
            ('Effective Gross Income', financial_metrics.get('effective_gross_income', 0))
        ]
        
        for item, value in income_data:
            ws[f'A{row}'] = item
            ws[f'B{row}'] = f"${value:,.0f}" if value else "$0"
            row += 1
        
        # Expense analysis
        row += 1
        ws[f'A{row}'] = "Expense Analysis"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        # Add more detailed financial breakdowns here
        # This would include operating expenses, debt service, etc.
        
    def _create_investment_projections_sheet(self, wb: Workbook, analysis_data: Dict[str, Any], include_charts: bool = True):
        """Create investment projections sheet with 5-year outlook"""
        
        ws = wb.create_sheet("Investment Projections")
        
        # Title
        ws['A1'] = "5-Year Investment Projections"
        ws['A1'].font = Font(size=16, bold=True, color=self.brand_colors['primary'])
        ws.merge_cells('A1:H1')
        
        # Create 5-year projection table
        row = 3
        
        # Headers
        headers = ['Year', 'Gross Income', 'Operating Expenses', 'NOI', 'Debt Service', 'Cash Flow', 'Cumulative Cash Flow']
        for col, header in enumerate(headers, 1):
            ws.cell(row=row, column=col, value=header).font = Font(bold=True)
        
        row += 1
        
        # Generate projections (simplified example)
        base_income = analysis_data.get('financial_metrics', {}).get('gross_income', 100000)
        base_expenses = base_income * 0.4  # 40% expense ratio
        base_debt_service = analysis_data.get('financial_metrics', {}).get('debt_service', 0)
        
        cumulative_cash_flow = 0
        
        for year in range(1, 6):
            # Apply growth rates
            projected_income = base_income * (1.03 ** year)  # 3% annual growth
            projected_expenses = base_expenses * (1.025 ** year)  # 2.5% annual growth
            projected_noi = projected_income - projected_expenses
            projected_cash_flow = projected_noi - base_debt_service
            cumulative_cash_flow += projected_cash_flow
            
            # Add row
            ws.cell(row=row, column=1, value=year)
            ws.cell(row=row, column=2, value=f"${projected_income:,.0f}")
            ws.cell(row=row, column=3, value=f"${projected_expenses:,.0f}")
            ws.cell(row=row, column=4, value=f"${projected_noi:,.0f}")
            ws.cell(row=row, column=5, value=f"${base_debt_service:,.0f}")
            ws.cell(row=row, column=6, value=f"${projected_cash_flow:,.0f}")
            ws.cell(row=row, column=7, value=f"${cumulative_cash_flow:,.0f}")
            
            row += 1
        
        # Add charts if requested
        if include_charts:
            self._add_cash_flow_chart(ws, row + 2)
    
    def _create_risk_assessment_sheet(self, wb: Workbook, analysis_data: Dict[str, Any]):
        """Create risk assessment sheet"""
        
        ws = wb.create_sheet("Risk Assessment")
        
        # Title
        ws['A1'] = "Investment Risk Assessment"
        ws['A1'].font = Font(size=16, bold=True, color=self.brand_colors['primary'])
        ws.merge_cells('A1:F1')
        
        # Risk factors
        row = 3
        ws[f'A{row}'] = "Identified Risk Factors"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        risk_factors = analysis_data.get('risk_factors', [])
        
        if risk_factors:
            for i, risk in enumerate(risk_factors, 1):
                ws[f'A{row}'] = f"{i}. {risk}"
                row += 1
        else:
            ws[f'A{row}'] = "No significant risk factors identified."
        
        # Risk mitigation strategies
        row += 2
        ws[f'A{row}'] = "Risk Mitigation Strategies"
        ws[f'A{row}'].font = Font(size=14, bold=True, color=self.brand_colors['primary'])
        row += 1
        
        # Add generic mitigation strategies
        mitigation_strategies = [
            "Conduct thorough due diligence and property inspection",
            "Obtain comprehensive property insurance coverage",
            "Maintain adequate cash reserves for unexpected expenses",
            "Consider professional property management services",
            "Monitor local market conditions and regulatory changes"
        ]
        
        for i, strategy in enumerate(mitigation_strategies, 1):
            ws[f'A{row}'] = f"{i}. {strategy}"
            row += 1
    
    def _create_property_details_sheet(self, wb: Workbook, property_data: Dict[str, Any]):
        """Create detailed property information sheet"""
        
        ws = wb.create_sheet("Property Details")
        
        # Title
        ws['A1'] = "Detailed Property Information"
        ws['A1'].font = Font(size=16, bold=True, color=self.brand_colors['primary'])
        ws.merge_cells('A1:F1')
        
        row = 3
        
        # All property details
        property_fields = [
            ('Property ID', property_data.get('id')),
            ('Address', property_data.get('address')),
            ('City', property_data.get('city')),
            ('State', property_data.get('state')),
            ('ZIP Code', property_data.get('zip_code')),
            ('Property Type', property_data.get('property_type')),
            ('Year Built', property_data.get('year_built')),
            ('Total Units', property_data.get('units')),
            ('Total Square Feet', property_data.get('total_sqft')),
            ('Lot Size', property_data.get('lot_size')),
            ('Foreclosure Stage', property_data.get('foreclosure_stage')),
            ('Outstanding Debt', f"${property_data.get('outstanding_debt', 0):,.0f}"),
            ('Estimated Value', f"${property_data.get('estimated_value', 0):,.0f}"),
            ('Minimum Bid', f"${property_data.get('minimum_bid', 0):,.0f}"),
            ('Auction Date', str(property_data.get('auction_date', ''))[:10] if property_data.get('auction_date') else 'N/A'),
            ('Discovery Date', str(property_data.get('discovery_date', ''))[:10] if property_data.get('discovery_date') else 'N/A'),
            ('Last Modified', str(property_data.get('last_modified', ''))[:10] if property_data.get('last_modified') else 'N/A')
        ]
        
        for field, value in property_fields:
            ws[f'A{row}'] = field
            ws[f'B{row}'] = value if value is not None else 'N/A'
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
        
        # Auto-adjust column widths
        for col in ['A', 'B']:
            ws.column_dimensions[col].auto_size = True
    
    def _create_raw_data_sheet(self, wb: Workbook, property_data: Dict[str, Any], analysis_data: Dict[str, Any]):
        """Create raw data sheet"""
        
        ws = wb.create_sheet("Raw Data")
        
        ws['A1'] = "Raw Analysis Data (JSON Format)"
        ws['A1'].font = Font(size=14, bold=True)
        
        # Convert data to JSON string for display
        raw_data = {
            'property_data': property_data,
            'analysis_data': analysis_data
        }
        
        json_str = json.dumps(raw_data, indent=2, default=str)
        
        # Split into lines and add to sheet
        lines = json_str.split('\n')
        for i, line in enumerate(lines[:1000], 3):  # Limit to first 1000 lines
            ws[f'A{i}'] = line
    
    def _add_cash_flow_chart(self, ws, start_row: int):
        """Add cash flow chart to worksheet"""
        
        try:
            # Create line chart
            chart = LineChart()
            chart.title = "5-Year Cash Flow Projection"
            chart.style = 10
            chart.y_axis.title = 'Cash Flow ($)'
            chart.x_axis.title = 'Year'
            
            # Data range (assuming cash flow is in column F, rows 4-8)
            data = Reference(ws, min_col=6, min_row=3, max_row=8, max_col=6)
            cats = Reference(ws, min_col=1, min_row=4, max_row=8)
            
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            
            # Add chart to worksheet
            ws.add_chart(chart, f"A{start_row}")
            
        except Exception as e:
            self.logger.warning(f"Could not add chart: {e}")
    
    # Utility Methods
    
    def _extract_property_basic_info(self, property_data: Dict[str, Any]) -> Dict[str, Any]:
        """Extract basic property information"""
        return {
            'id': property_data.get('id'),
            'address': property_data.get('address'),
            'city': property_data.get('city'),
            'state': property_data.get('state'),
            'property_type': property_data.get('property_type'),
            'units': property_data.get('units'),
            'total_sqft': property_data.get('total_sqft')
        }
    
    def _extract_foreclosure_details(self, property_data: Dict[str, Any]) -> Dict[str, Any]:
        """Extract foreclosure-specific details"""
        return {
            'stage': property_data.get('foreclosure_stage'),
            'outstanding_debt': property_data.get('outstanding_debt'),
            'estimated_value': property_data.get('estimated_value'),
            'minimum_bid': property_data.get('minimum_bid'),
            'auction_date': property_data.get('auction_date'),
            'case_number': property_data.get('case_number')
        }
    
    def _extract_location_info(self, property_data: Dict[str, Any]) -> Dict[str, Any]:
        """Extract location information"""
        return {
            'address': property_data.get('address'),
            'city': property_data.get('city'),
            'state': property_data.get('state'),
            'zip_code': property_data.get('zip_code'),
            'county': property_data.get('county'),
            'coordinates': {
                'latitude': property_data.get('latitude'),
                'longitude': property_data.get('longitude')
            }
        }
    
    def _extract_analysis_summary(self, analysis_data: Dict[str, Any]) -> Dict[str, Any]:
        """Extract analysis summary"""
        return {
            'investment_grade': analysis_data.get('investment_grade'),
            'viability_score': analysis_data.get('viability_score'),
            'confidence_score': analysis_data.get('confidence_score'),
            'analysis_date': analysis_data.get('analysis_date'),
            'key_strengths': analysis_data.get('ai_insights', {}).get('key_strengths', []),
            'key_concerns': analysis_data.get('ai_insights', {}).get('key_concerns', [])
        }
    
    def _generate_export_summary(self, export_session: Dict[str, Any], export_path: str) -> Dict[str, Any]:
        """Generate export summary file"""
        
        try:
            summary = {
                'export_summary': {
                    'export_id': export_session['export_id'],
                    'created_at': export_session['created_at'].isoformat(),
                    'completed_at': export_session.get('completed_at', datetime.now()).isoformat(),
                    'property_id': export_session['property_id'],
                    'analysis_id': export_session['analysis_id'],
                    'requested_formats': export_session['requested_formats'],
                    'generated_files': len(export_session['files']),
                    'total_size_mb': sum(f.get('size_mb', 0) for f in export_session['files'])
                },
                'files': export_session['files']
            }
            
            filename = 'export_summary.json'
            file_path = os.path.join(export_path, filename)
            
            with open(file_path, 'w') as f:
                json.dump(summary, f, indent=2, default=str)
            
            return {
                'type': 'summary',
                'filename': filename,
                'path': file_path,
                'size_mb': os.path.getsize(file_path) / (1024 * 1024),
                'description': 'Export session summary and file inventory'
            }
            
        except Exception as e:
            self.logger.error(f"Error generating export summary: {str(e)}")
            return None
    
    async def cleanup_old_exports(self, days_threshold: int = None) -> Dict[str, Any]:
        """Clean up old export files"""
        
        try:
            if days_threshold is None:
                days_threshold = self.max_export_age_days
            
            cutoff_date = datetime.now() - timedelta(days=days_threshold)
            
            cleanup_stats = {
                'exports_cleaned': 0,
                'files_deleted': 0,
                'space_freed_mb': 0,
                'errors': []
            }
            
            if not os.path.exists(self.export_dir):
                return cleanup_stats
            
            # Clean up old export directories
            for item in os.listdir(self.export_dir):
                item_path = os.path.join(self.export_dir, item)
                
                if os.path.isdir(item_path):
                    try:
                        # Check directory age
                        dir_mtime = datetime.fromtimestamp(os.path.getmtime(item_path))
                        
                        if dir_mtime < cutoff_date:
                            # Calculate size before deletion
                            dir_size = 0
                            file_count = 0
                            
                            for root, dirs, files in os.walk(item_path):
                                for file in files:
                                    file_path = os.path.join(root, file)
                                    dir_size += os.path.getsize(file_path)
                                    file_count += 1
                            
                            # Delete directory
                            shutil.rmtree(item_path)
                            
                            cleanup_stats['exports_cleaned'] += 1
                            cleanup_stats['files_deleted'] += file_count
                            cleanup_stats['space_freed_mb'] += dir_size / (1024 * 1024)
                            
                    except Exception as e:
                        cleanup_stats['errors'].append(f"Error cleaning {item}: {str(e)}")
            
            self.logger.info(f"Export cleanup completed: {cleanup_stats}")
            return cleanup_stats
            
        except Exception as e:
            self.logger.error(f"Error during export cleanup: {str(e)}")
            raise
    
    # Placeholder methods for portfolio exports (simplified implementations)
    
    async def _export_portfolio_json(self, portfolio_summary, properties, export_path, include_details):
        """Export portfolio data as JSON"""
        filename = f"portfolio_summary_{datetime.now().strftime('%Y%m%d')}.json"
        file_path = os.path.join(export_path, filename)
        
        data = {
            'portfolio_summary': portfolio_summary,
            'properties': properties if include_details else [{'id': p['id'], 'address': p.get('address', '')} for p in properties]
        }
        
        with open(file_path, 'w') as f:
            json.dump(data, f, indent=2, default=str)
        
        return {
            'type': 'json',
            'filename': filename,
            'path': file_path,
            'size_mb': os.path.getsize(file_path) / (1024 * 1024),
            'description': 'Portfolio summary data'
        }
    
    async def _export_portfolio_excel(self, portfolio_summary, properties, export_path, include_details):
        """Export portfolio as Excel file"""
        if not EXCEL_AVAILABLE:
            return None
        
        filename = f"portfolio_analysis_{datetime.now().strftime('%Y%m%d')}.xlsx"
        file_path = os.path.join(export_path, filename)
        
        wb = Workbook()
        wb.remove(wb.active)
        
        # Summary sheet
        ws = wb.create_sheet("Portfolio Summary")
        ws['A1'] = "Portfolio Analysis Summary"
        ws['A1'].font = Font(size=16, bold=True)
        
        # Properties list sheet
        if include_details:
            ws2 = wb.create_sheet("Property List")
            ws2['A1'] = "Property List"
            
            # Add properties data (simplified)
            headers = ['ID', 'Address', 'City', 'State', 'Stage', 'Grade', 'Investment Score']
            for col, header in enumerate(headers, 1):
                ws2.cell(row=2, column=col, value=header).font = Font(bold=True)
            
            for row, prop in enumerate(properties[:100], 3):  # Limit to first 100
                ws2.cell(row=row, column=1, value=prop.get('id', ''))
                ws2.cell(row=row, column=2, value=prop.get('address', ''))
                ws2.cell(row=row, column=3, value=prop.get('city', ''))
                ws2.cell(row=row, column=4, value=prop.get('state', ''))
                ws2.cell(row=row, column=5, value=prop.get('foreclosure_stage', ''))
                ws2.cell(row=row, column=6, value=prop.get('investment_grade', ''))
                ws2.cell(row=row, column=7, value=prop.get('investment_score', 0))
        
        wb.save(file_path)
        
        return {
            'type': 'excel',
            'filename': filename,
            'path': file_path,
            'size_mb': os.path.getsize(file_path) / (1024 * 1024),
            'description': 'Portfolio analysis spreadsheet'
        }
    
    async def _export_portfolio_pdf(self, portfolio_summary, properties, export_path):
        """Export portfolio as PDF report"""
        if not PDF_AVAILABLE:
            return None
        
        filename = f"portfolio_report_{datetime.now().strftime('%Y%m%d')}.pdf"
        file_path = os.path.join(export_path, filename)
        
        doc = SimpleDocTemplate(file_path, pagesize=letter)
        styles = getSampleStyleSheet()
        content = []
        
        # Title
        title = Paragraph("Portfolio Analysis Report", styles['Title'])
        content.append(title)
        content.append(Spacer(1, 20))
        
        # Summary statistics
        summary_text = f"""
        <b>Portfolio Overview:</b><br/>
        Total Properties: {portfolio_summary['summary']['total_properties']}<br/>
        Average Investment Score: {portfolio_summary['financial_metrics'].get('average_investment_score', 0):.2f}<br/>
        Total Estimated Value: ${portfolio_summary['financial_metrics'].get('total_estimated_value', 0):,.0f}
        """
        
        content.append(Paragraph(summary_text, styles['Normal']))
        
        doc.build(content)
        
        return {
            'type': 'pdf',
            'filename': filename,
            'path': file_path,
            'size_mb': os.path.getsize(file_path) / (1024 * 1024),
            'description': 'Portfolio analysis PDF report'
        }
    
    # Opportunity export methods (simplified)
    
    async def _export_opportunities_pdf(self, opportunities, export_path):
        """Export opportunities as PDF"""
        # Implementation similar to portfolio PDF but focused on opportunities
        filename = f"investment_opportunities_{datetime.now().strftime('%Y%m%d')}.pdf"
        file_path = os.path.join(export_path, filename)
        
        # Simplified PDF creation
        with open(file_path, 'w') as f:
            f.write("Investment Opportunities Report - PDF generation requires reportlab library")
        
        return {
            'type': 'pdf',
            'filename': filename,
            'path': file_path,
            'size_mb': os.path.getsize(file_path) / (1024 * 1024),
            'description': 'Investment opportunities PDF report'
        }
    
    async def _export_opportunities_excel(self, opportunities, export_path):
        """Export opportunities as Excel"""
        # Implementation similar to portfolio Excel but focused on opportunities
        return None  # Placeholder
    
    async def _export_opportunities_json(self, opportunities, export_path):
        """Export opportunities as JSON"""
        filename = f"investment_opportunities_{datetime.now().strftime('%Y%m%d')}.json"
        file_path = os.path.join(export_path, filename)
        
        with open(file_path, 'w') as f:
            json.dump({
                'opportunities': opportunities,
                'generated_at': datetime.now().isoformat(),
                'total_count': len(opportunities)
            }, f, indent=2, default=str)
        
        return {
            'type': 'json',
            'filename': filename,
            'path': file_path,
            'size_mb': os.path.getsize(file_path) / (1024 * 1024),
            'description': 'Investment opportunities data'
        }